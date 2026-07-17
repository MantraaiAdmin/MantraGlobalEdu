#!/usr/bin/env python3
"""
Mantra Global Education — Student Webinar PowerPoint (Sales Edition)
Senior sales narrative: problem → proof → process → partnership → close
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ASSETS = Path(__file__).parent / "ppt-assets"
OUTPUT = Path(__file__).parent / "Mantra-Student-Webinar.pptx"

PRIMARY = RGBColor(0, 35, 78)
ACCENT = RGBColor(200, 145, 22)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(90, 107, 125)
LIGHT = RGBColor(247, 249, 252)
SOFT_GOLD = RGBColor(245, 234, 212)
RED = RGBColor(192, 57, 43)
GREEN = RGBColor(39, 174, 96)
GREY_TXT = RGBColor(160, 160, 160)


# ── Helpers ──────────────────────────────────────────────────────────────

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.07))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()


def rect(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=12, bold=False, color=PRIMARY,
        align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.name = font
    p.font.color.rgb = color
    p.alignment = align
    return box


def header(slide, label, title, num, subtitle=None, dark=False):
    tc = WHITE if dark else PRIMARY
    mc = GREY_TXT if dark else MUTED
    txt(slide, Inches(0.55), Inches(0.32), Inches(5), Inches(0.28), label, 8, True, ACCENT)
    txt(slide, Inches(0.55), Inches(0.55), Inches(10), Inches(0.5), title, 26, True, tc, font="Georgia")
    if subtitle:
        txt(slide, Inches(0.55), Inches(1.05), Inches(11), Inches(0.35), subtitle, 12, False, mc)
    txt(slide, Inches(11.7), Inches(0.38), Inches(1.3), Inches(0.3), num, 10, False, mc, PP_ALIGN.RIGHT)


def footer(slide, left="Mantra Global Education", right="mantraglobaledu.com"):
    txt(slide, Inches(0.55), Inches(7.1), Inches(5), Inches(0.25), left, 8, False, MUTED)
    txt(slide, Inches(8), Inches(7.1), Inches(4.8), Inches(0.25), right, 8, False, MUTED, PP_ALIGN.RIGHT)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def bullets(slide, l, t, w, h, items, sz=11, color=MUTED, bullet_color=ACCENT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.level = 0
        p.space_after = Pt(6)
    return box


def metric_card(slide, l, t, w, h, value, label, sub=None):
    box = rect(slide, l, t, w, h, PRIMARY)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(8)
    p2.font.color.rgb = GREY_TXT
    p2.alignment = PP_ALIGN.CENTER
    if sub:
        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(7)
        p3.font.color.rgb = GREY_TXT
        p3.alignment = PP_ALIGN.CENTER


def pic(slide, name, l, t, w, h):
    p = ASSETS / name
    if p.exists() and p.stat().st_size > 500:
        slide.shapes.add_picture(str(p), l, t, w, h)


def section_card(slide, l, t, w, h, title, body, accent_bar=True):
    rect(slide, l, t, w, h, LIGHT)
    if accent_bar:
        rect(slide, l, t, Inches(0.06), h, ACCENT)
    txt(slide, l + Inches(0.18), t + Inches(0.12), w - Inches(0.3), Inches(0.35),
        title, 12, True, PRIMARY)
    txt(slide, l + Inches(0.18), t + Inches(0.48), w - Inches(0.3), h - Inches(0.55),
        body, 10, False, MUTED)


# ── Build ────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    N = "14"

    # ═══ SLIDE 1 — COVER ═══
    s = prs.slides.add_slide(B)
    bg(s, PRIMARY)
    bar(s)
    pic(s, "campus.jpg", Inches(7.2), Inches(0.07), Inches(6.13), Inches(7.43))
    overlay = rect(s, Inches(7.2), Inches(0.07), Inches(6.13), Inches(7.43), PRIMARY)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = PRIMARY
    # semi-transparent effect via second overlay strip
    txt(s, Inches(0.7), Inches(0.55), Inches(3), Inches(0.3),
        "IN-PERSON STUDENT SESSION", 9, True, ACCENT)
    txt(s, Inches(0.7), Inches(0.9), Inches(6.5), Inches(1.1),
        "Mantra Global\nEducation", 40, True, WHITE, font="Georgia")
    txt(s, Inches(0.7), Inches(2.2), Inches(6), Inches(0.35),
        "Your Journey. Our Guidance. Global Success.", 13, False, ACCENT)
    txt(s, Inches(0.7), Inches(2.85), Inches(6.2), Inches(1.3),
        "From Your Campus to the World's Best Universities\n"
        "USA  ·  United Kingdom  ·  Australia",
        18, True, WHITE, font="Georgia")
    txt(s, Inches(0.7), Inches(4.35), Inches(6), Inches(1.0),
        "A 45-minute session to help you choose the right country, "
        "university, and pathway — with a certified counselor who has "
        "guided 700+ students to global success.",
        12, False, GREY_TXT)
    rect(s, Inches(0.7), Inches(5.6), Inches(6), Inches(0.9), RGBColor(13, 50, 90))
    txt(s, Inches(0.9), Inches(5.72), Inches(5.6), Inches(0.7),
        "Today's outcome: You leave with a clear study-abroad action plan "
        "and a free 1-on-1 counseling slot booked.",
        11, True, SOFT_GOLD)
    txt(s, Inches(0.7), Inches(6.65), Inches(5), Inches(0.3),
        "mantraglobaledu.com  |  HQ: Coimbatore, Tamil Nadu", 10, False, GREY_TXT)
    notes(s, "OPEN WITH ENERGY. 'Raise your hand if you've thought about studying abroad but don't know where to start.' "
             "Promise: they leave with clarity + free counseling booked.")
    footer(s, right="Slide 01 / " + N)

    # ═══ SLIDE 2 — AGENDA (RICH) ═══
    s = prs.slides.add_slide(B)
    bg(s, WHITE)
    bar(s)
    header(s, "SESSION ROADMAP", "What We Will Cover — And What You Will Gain",
         "02 / " + N,
         "Every section is designed to answer a real question you have right now.")
    agenda = [
        ("01", "5 min", "Know Your Starting Point",
         "Interactive poll — where do you want to study and why?",
         "You'll identify your top destination priority"),
        ("02", "8 min", "The Global Opportunity",
         "Career ROI, salary data, and why timing matters for 2026 intake",
         "You'll understand the financial and career case for going abroad"),
        ("03", "7 min", "Bust the Myths",
         "Cost, IELTS, DIY vs. agent — separating fact from fear",
         "You'll stop believing the 4 myths that block most students"),
        ("04", "5 min", "Real Student Story",
         "Video — a journey from campus to international university",
         "You'll see proof this is achievable for students like you"),
        ("05", "12 min", "Country Masterclass",
         "USA, UK, Australia — programs, visas, scholarships, timelines",
         "You'll know which country fits YOUR profile and goals"),
        ("06", "8 min", "The MANTRA Method",
         "Our 6-stage framework — from counseling to career abroad",
         "You'll see exactly how we take you from confused to enrolled"),
    ]
    y0 = Inches(1.45)
    for i, (num, time, title, desc, gain) in enumerate(agenda):
        y = y0 + Inches(i * 0.88)
        rect(s, Inches(0.55), y, Inches(0.45), Inches(0.72), PRIMARY)
        tf = s.shapes[-1].text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = ACCENT
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        txt(s, Inches(1.1), y + Inches(0.02), Inches(0.7), Inches(0.25), time, 8, True, ACCENT)
        txt(s, Inches(1.1), y + Inches(0.22), Inches(5.5), Inches(0.28), title, 12, True, PRIMARY)
        txt(s, Inches(1.1), y + Inches(0.48), Inches(5.8), Inches(0.22), desc, 9, False, MUTED)
    # Right panel — outcomes
    panel = rect(s, Inches(7.2), Inches(1.45), Inches(5.6), Inches(5.3), PRIMARY)
    txt(s, Inches(7.45), Inches(1.6), Inches(5.1), Inches(0.35),
        "BY THE END OF THIS SESSION", 10, True, ACCENT)
    txt(s, Inches(7.45), Inches(1.95), Inches(5.1), Inches(0.5),
        "You Will Walk Away Knowing:", 16, True, WHITE, font="Georgia")
    outcomes = [
        "Which country (USA / UK / AUS) best matches your career goal",
        "Realistic budget, scholarship, and education loan options",
        "Exact intake timeline — when to apply, test, and prepare",
        "Why Mantra's 1-on-1 model beats generic consultancies",
        "Your personal next step — with a free counseling slot booked",
    ]
    bullets(s, Inches(7.45), Inches(2.55), Inches(5.1), Inches(3.5), outcomes, 11, GREY_TXT)
    rect(s, Inches(7.45), Inches(5.9), Inches(5.1), Inches(0.65), RGBColor(13, 50, 90))
    txt(s, Inches(7.6), Inches(6.0), Inches(4.8), Inches(0.45),
        "Sections 07–08: Why Mantra + Q&A & Free Counseling Booking",
        10, True, SOFT_GOLD, PP_ALIGN.CENTER)
    notes(s, "Don't rush this slide. Students need to see VALUE of staying 45 min. "
             "Point to right panel: 'These are YOUR takeaways.'")
    footer(s)

    # ═══ SLIDE 3 — PAIN POINTS ═══
    s = prs.slides.add_slide(B)
    bg(s, WHITE)
    bar(s)
    header(s, "THE REALITY CHECK", "Why Most Students Never Make It Abroad",
         "03 / " + N,
         "It's not lack of ambition. It's lack of the right guidance at the right time.")
    pains = [
        ("Confusion at Final Year",
         "60% of engineering students want to study abroad but haven't started "
         "research by 7th semester. By the time they act, top intake deadlines have passed.",
         "Start 12–15 months before your target intake"),
        ("Wrong University Choices",
         "Students apply to universities based on ads, not profile fit. "
         "Result: rejections, wasted application fees, and lost confidence.",
         "Profile-matched shortlisting saves time and money"),
        ("Hidden Costs & Surprises",
         "Many agents quote low fees upfront, then add charges for SOP, visa, "
         "and documentation. Families lose trust and students drop the plan.",
         "Mantra: transparent process, no hidden surprises"),
        ("Going Alone on Applications",
         "SOP, LOR, visa documentation, and scholarship essays require strategy. "
         "One weak essay can cost an admission offer worth ₹40–80 lakhs in career value.",
         "Expert review on every document before submission"),
    ]
    for i, (title, body, fix) in enumerate(pains):
        row, col = divmod(i, 2)
        x = Inches(0.55 + col * 6.4)
        y = Inches(1.5 + row * 2.65)
        section_card(s, x, y, Inches(6.1), Inches(2.4), title, body)
        rect(s, x + Inches(0.18), y + Inches(1.85), Inches(5.7), Inches(0.42), SOFT_GOLD)
        txt(s, x + Inches(0.28), y + Inches(1.92), Inches(5.5), Inches(0.3),
            f"Mantra fix: {fix}", 9, True, PRIMARY)
    notes(s, "AGITATE THE PAIN. Ask: 'Who feels behind on applications?' "
             "This positions Mantra as the solution before you even introduce the framework.")
    footer(s)

    # ═══ SLIDE 4 — GLOBAL OPPORTUNITY ═══
    s = prs.slides.add_slide(B)
    bg(s, WHITE)
    bar(s)
    header(s, "THE BUSINESS CASE", "Why a Global Postgraduate Degree Changes Everything",
         "04 / " + N,
         "This is not about leaving India. It's about multiplying your career options.")
    pic(s, "graduates.jpg", Inches(0.55), Inches(1.45), Inches(4.8), Inches(3.8))
    metrics = [
        ("$95K–$140K", "Avg. US tech salary\npost-MS (Year 3)", "3–5x India equivalent"),
        ("12 months", "UK Masters duration", "Enter workforce faster"),
        ("4 years", "Australia PSW visa", "Work right after graduation"),
        ("40+", "PG programs we cover", "CS, Data, MBA, Engineering"),
    ]
    for i, (v, l, sub) in enumerate(metrics):
        metric_card(s, Inches(5.6 + (i % 2) * 3.85), Inches(1.45 + (i // 2) * 1.55),
                    Inches(3.6), Inches(1.35), v, l, sub)
    section_card(s, Inches(5.6), Inches(4.65), Inches(7.2), Inches(2.1),
                 "What Top Employers Look For",
                 "Global exposure · International degree · Cross-cultural communication · "
                 "Specialized skills from world-ranked universities · "
                 "Post-study work experience in USA, UK, or Australia\n\n"
                 "An international PG is not an expense — it is the highest-ROI investment "
                 "in your 20s. The question is not IF you can afford it, but HOW to plan it smartly.")
    notes(s, "Use numbers confidently. Engineers care about ROI. "
             "Say: 'A 2-year US MS can pay for itself within 3 years of working there.'")
    footer(s)

    # ═══ SLIDE 5 — ICEBREAKER ═══
    s = prs.slides.add_slide(B)
    bg(s, WHITE)
    bar(s)
    header(s, "INTERACTIVE · LET'S START", "Where Do You See Yourself in 2027?",
         "05 / " + N,
         "Raise your hand or call out — there is no wrong answer. This helps me tailor the session.")
    banner = rect(s, Inches(0.55), Inches(1.4), Inches(12.2), Inches(0.65), PRIMARY)
    tf = banner.text_frame
    tf.paragraphs[0].text = (
        "QUESTION: If you could study abroad tomorrow, which country would you choose — "
        "USA, UK, or Australia? And what do you want to study?"
    )
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    countries = [
        ("usa.jpg", "UNITED STATES",
         "Best for: MS Computer Science, Data Science, AI, MBA\n"
         "Strength: Highest salary ceiling + 3-year STEM OPT\n"
         "Intake: Fall (Aug) & Spring (Jan) | Tests: GRE + IELTS/TOEFL\n"
         "Ideal student: Wants maximum career ROI & tech leadership roles"),
        ("uk.jpg", "UNITED KINGDOM",
         "Best for: 1-year MSc, Business, Engineering, Law\n"
         "Strength: Fastest path — graduate in 12 months, work 2 years\n"
         "Intake: September & January | Tests: IELTS (GRE often waived)\n"
         "Ideal student: Wants quality degree quickly with lower total cost"),
        ("australia.jpg", "AUSTRALIA",
         "Best for: IT, Engineering, Healthcare, Business Analytics\n"
         "Strength: Post-study work visa up to 4 years + PR pathway\n"
         "Intake: February & July | Tests: IELTS/PTE\n"
         "Ideal student: Wants work-life balance + long-term settlement options"),
    ]
    for i, (img, name, detail) in enumerate(countries):
        x = Inches(0.55 + i * 4.25)
        pic(s, img, x, Inches(2.2), Inches(4.0), Inches(1.85))
        rect(s, x, Inches(4.1), Inches(4.0), Inches(2.65), LIGHT)
        txt(s, x + Inches(0.12), Inches(4.18), Inches(3.7), Inches(0.3), name, 12, True, PRIMARY)
        txt(s, x + Inches(0.12), Inches(4.5), Inches(3.7), Inches(2.1), detail, 9, False, MUTED)
    notes(s, "PAUSE 3 minutes. Count hands. Mirror back: 'I see many of you leaning USA — "
             "we'll deep-dive that in a few minutes.' Makes it personal.")
    footer(s)

    # ═══ SLIDE 6 — MYTHS ═══
    s = prs.slides.add_slide(B)
    bg(s, WHITE)
    bar(s)
    header(s, "OBJECTION HANDLING", "4 Myths That Stop Students — And The Truth",
         "06 / " + N,
         "Every year, talented students delay or drop their plans because of these misconceptions.")
    myths = [
        ("MYTH 1", "Only rich families can afford study abroad",
         "REALITY",
         "Total cost for UK MSc: ₹25–35L (1 year). US MS: ₹40–60L (2 years). "
         "Scholarships cover 20–100% at many universities. Education loans from "
         "SBI, HDFC Credila, Prodigy Finance cover 80–100% with moratorium until graduation. "
         "We help families plan a realistic budget — most middle-class families CAN make this work."),
        ("MYTH 2", "You need IELTS 7.5+ or perfect scores",
         "REALITY",
         "Many universities accept IELTS 6.0–6.5. Some offer conditional admission "
         "(pre-sessional English). USA accepts TOEFL/Duolingo at select unis. "
         "We match you to universities where YOUR score is competitive — not where it isn't."),
        ("MYTH 3", "I can apply myself — why pay a consultant?",
         "REALITY",
         "Application is 20% of the journey. The other 80%: university shortlisting, "
         "SOP strategy, scholarship negotiation, visa documentation, financial planning, "
         "pre-departure, and post-arrival support. One rejected visa = 1 year lost. "
         "Our students get it right the first time because of expert review at every step."),
        ("MYTH 4", "All consultancies do the same thing",
         "REALITY",
         "Call-center agents handle 50 students each. Mantra assigns ONE certified counselor "
         "to YOU from day one. We focus on only 3 countries (not 30). We have 15+ direct "
         "university partnerships. 700+ placements. British Council certified. "
         "The difference is personal — and it shows in admission results."),
    ]
    for i, (ml, mt, tl, tt) in enumerate(myths):
        row, col = divmod(i, 2)
        x = Inches(0.55 + col * 6.4)
        y = Inches(1.45 + row * 2.85)
        rect(s, x, y, Inches(6.1), Inches(0.9), RGBColor(253, 232, 232))
        txt(s, x + Inches(0.12), y + Inches(0.06), Inches(0.8), Inches(0.2), ml, 7, True, RED)
        txt(s, x + Inches(0.12), y + Inches(0.28), Inches(5.8), Inches(0.55), mt, 11, True, PRIMARY)
        rect(s, x, y + Inches(0.95), Inches(6.1), Inches(1.7), RGBColor(232, 245, 233))
        txt(s, x + Inches(0.12), y + Inches(1.02), Inches(0.8), Inches(0.2), tl, 7, True, GREEN)
        txt(s, x + Inches(0.12), y + Inches(1.24), Inches(5.8), Inches(1.3), tt, 9, False, PRIMARY)
    notes(s, "Ask before each myth: 'Who has heard this?' Validates their concern, then you solve it.")
    footer(s)

    # ═══ SLIDE 7 — VIDEO ═══
    s = prs.slides.add_slide(B)
    bg(s, PRIMARY)
    bar(s)
    header(s, "PROOF, NOT PROMISES", "Watch: A Student Journey With Mantra",
         "07 / " + N,
         "Real students. Real universities. Real outcomes.", dark=True)
    vp = ASSETS / "student-journey.mp4"
    poster = ASSETS / "graduates.jpg"
    if vp.exists():
        try:
            kw = {"mime_type": "video/mp4"}
            if poster.exists() and poster.stat().st_size > 500:
                kw["poster_frame_image"] = str(poster)
            s.shapes.add_movie(str(vp), Inches(1.2), Inches(1.5), Inches(8.5), Inches(4.8), **kw)
        except Exception:
            pic(s, "graduates.jpg", Inches(1.2), Inches(1.5), Inches(8.5), Inches(4.8))
    section_card(s, Inches(9.9), Inches(1.5), Inches(3.0), Inches(4.8),
                 "After the video, discuss:",
                 "• What surprised you about the process?\n"
                 "• Which step seemed hardest?\n"
                 "• Could you see yourself in this journey?\n\n"
                 "This is the same structured path we build for every Mantra student — "
                 "personalized to your profile, budget, and career goal.")
    txt(s, Inches(1.2), Inches(6.45), Inches(11), Inches(0.3),
        "Click video to play  |  ~2 minutes  |  Then open floor for 2 student reactions",
        10, False, GREY_TXT, PP_ALIGN.CENTER)
    notes(s, "After video: pick 2 students to share one takeaway. Social proof is strongest from peers.")
    footer(s, right="Slide 07 / " + N)

    # ═══ SLIDE 8 — USA DEEP DIVE ═══
    s = prs.slides.add_slide(B)
    bg(s, WHITE)
    bar(s)
    header(s, "COUNTRY MASTERCLASS", "United States — The Career Accelerator",
         "08 / " + N,
         "Highest ROI for engineering & tech graduates. Most competitive, most rewarding.")
    pic(s, "usa.jpg", Inches(0.55), Inches(1.4), Inches(4.2), Inches(2.6))
    section_card(s, Inches(0.55), Inches(4.1), Inches(4.2), Inches(2.7),
                 "Who Should Choose USA?",
                 "Engineering & CS students targeting Silicon Valley, FAANG, or US tech roles. "
                 "Students comfortable with 2-year investment for long-term 3–5x salary growth. "
                 "GRE-ready students aiming for top-100 ranked universities.")
    cols = [
        ("Top Programs", "MS Computer Science\nMS Data Science & AI\nMS Electrical & Mechanical\nMBA (STEM-designated)"),
        ("Key Requirements", "GRE: 300+ (varies by uni)\nIELTS: 6.5+ or TOEFL 80+\nGPA: 3.0+ (60–70% Indian equiv.)\nSOP + 3 LORs + Resume"),
        ("Costs & Funding", "Tuition: $20K–$55K/year\nLiving: $12K–$18K/year\nScholarships: $5K–full ride\nLoans: Up to 100% coverage"),
        ("Post-Study Work", "STEM OPT: 3 years paid work\nH-1B sponsorship pathway\nAvg. starting: $85K–$120K\nROI break-even: ~3 years"),
    ]
    for i, (title, body) in enumerate(cols):
        x = Inches(5.0 + (i % 2) * 4.1)
        y = Inches(1.4 + (i // 2) * 2.55)
        section_card(s, x, y, Inches(3.85), Inches(2.35), title, body)
    rect(s, Inches(5.0), Inches(6.45), Inches(7.8), Inches(0.45), SOFT_GOLD)
    txt(s, Inches(5.15), Inches(6.52), Inches(7.5), Inches(0.3),
        "Mantra advantage: 15+ US partner universities with real fee data & admission requirements on file",
        9, True, PRIMARY, PP_ALIGN.CENTER)
    notes(s, "For engineers in the room: emphasize STEM OPT. Parents care about ROI — cite salary numbers.")
    footer(s)

    # ═══ SLIDE 9 — UK + AUS ═══
    s = prs.slides.add_slide(B)
    bg(s, WHITE)
    bar(s)
    header(s, "COUNTRY MASTERCLASS", "United Kingdom & Australia — Fast Track & Future Settlement",
         "09 / " + N,
         "Two powerful alternatives depending on your timeline, budget, and long-term plan.")
    for i, (img, country, tag, programs, reqs, cost, work) in enumerate([
        ("uk.jpg", "UNITED KINGDOM", "Best for speed & prestige",
         "1-year MSc (no wasted time)\nRussell Group universities\nMBA, Engineering, Finance, Law",
         "IELTS 6.0–6.5 (GRE often waived)\nStrong SOP + academic references\n60%+ in UG for most unis",
         "₹25–35L total (tuition + living)\nScholarships: £3K–£10K common\nChevening for leaders",
         "2-year Graduate Route visa\nWork full-time after graduation\nGateway to global careers in Europe"),
        ("australia.jpg", "AUSTRALIA", "Best for work rights & lifestyle",
         "MS IT, Engineering, Analytics\nHealthcare & Nursing pathways\nStrong PR pathway for skilled grads",
         "IELTS 6.0–6.5 or PTE 58+\nGPA: 60%+ for most programs\nWork experience helps for MBA",
         "₹30–45L total (2-year programs)\nAustralia Awards scholarships\nPart-time work: 48 hrs/fortnight",
         "Post-Study Work: 2–4 years\nPR points system favors graduates\nHigh demand for IT & engineering"),
    ]):
        x = Inches(0.55 + i * 6.4)
        pic(s, img, x, Inches(1.4), Inches(6.1), Inches(1.6))
        txt(s, x, Inches(3.05), Inches(6.1), Inches(0.35), country, 14, True, PRIMARY)
        txt(s, x, Inches(3.38), Inches(6.1), Inches(0.25), tag, 10, True, ACCENT)
        for j, (lbl, content) in enumerate([
            ("Programs", programs), ("Requirements", reqs),
            ("Investment", cost), ("After Graduation", work),
        ]):
            y = Inches(3.7 + j * 0.82)
            txt(s, x, y, Inches(1.2), Inches(0.25), lbl, 8, True, PRIMARY)
            txt(s, x + Inches(1.15), y, Inches(4.9), Inches(0.75), content, 9, False, MUTED)
    notes(s, "Compare directly: 'UK if you want speed. Australia if you want settlement. USA if you want max salary.'")
    footer(s)

    # ═══ SLIDE 10 — MANTRA FRAMEWORK (DETAILED) ═══
    s = prs.slides.add_slide(B)
    bg(s, WHITE)
    bar(s)
    header(s, "OUR PROPRIETARY METHOD", "The MANTRA Framework — Your End-to-End Success System",
         "10 / " + N,
         "Other agents stop at the application. We walk with you until you succeed abroad.")
    stages = [
        ("M", "MENTOR", "Discover Your Future",
         "1-on-1 career counseling, academic profile audit, budget planning, "
         "and goal-setting. We build your personalized study-abroad roadmap before any application begins.",
         "Outcome: Clear study plan with target intake"),
        ("A", "ASSESS", "Choose the Right University",
         "Data-driven shortlisting from 15+ partner universities. We match your GPA, "
         "test scores, budget, and career goal — not commission rates.",
         "Outcome: 5–8 best-fit universities identified"),
        ("N", "NAVIGATE", "Build a Winning Application",
         "Expert SOP writing, LOR guidance, resume enhancement, and document review. "
         "Every application is reviewed by counselors with 10+ years experience.",
         "Outcome: Competitive, error-free applications"),
        ("T", "TRANSITION", "Visa, Finance & Pre-Departure",
         "Visa documentation, mock interviews, education loan facilitation, "
         "scholarship applications, and financial planning for your family.",
         "Outcome: Visa approved, finances secured"),
        ("R", "RELOCATE", "Land & Settle Confidently",
         "Accommodation guidance, airport pickup coordination, forex, SIM, "
         "banking setup, and pre-departure orientation for students and parents.",
         "Outcome: Smooth arrival in host country"),
        ("A", "ARRIVE", "Thrive Beyond Admission",
         "Orientation support, alumni network access, and ongoing mentorship. "
         "Our relationship doesn't end at the airport — it begins there.",
         "Outcome: Academic success & career growth"),
    ]
    for i, (letter, word, title, desc, outcome) in enumerate(stages):
        row, col = divmod(i, 3)
        x = Inches(0.55 + col * 4.25)
        y = Inches(1.45 + row * 2.75)
        rect(s, x, y, Inches(4.0), Inches(2.5), LIGHT, PRIMARY)
        txt(s, x + Inches(0.12), y + Inches(0.1), Inches(0.5), Inches(0.45),
            letter, 22, True, ACCENT, font="Georgia")
        txt(s, x + Inches(0.55), y + Inches(0.12), Inches(3.3), Inches(0.25), word, 8, True, PRIMARY)
        txt(s, x + Inches(0.12), y + Inches(0.55), Inches(3.7), Inches(0.3), title, 11, True, PRIMARY)
        txt(s, x + Inches(0.12), y + Inches(0.88), Inches(3.7), Inches(1.1), desc, 8, False, MUTED)
        rect(s, x + Inches(0.12), y + Inches(2.05), Inches(3.7), Inches(0.35), SOFT_GOLD)
        txt(s, x + Inches(0.18), y + Inches(2.1), Inches(3.55), Inches(0.25), outcome, 8, True, PRIMARY)
    notes(s, "THIS IS YOUR DIFFERENTIATOR. Say slowly: 'We don't disappear after you apply. "
             "We are with you until you thrive.' Students and parents remember this.")
    footer(s)

    # ═══ SLIDE 11 — WHY MANTRA ═══
    s = prs.slides.add_slide(B)
    bg(s, WHITE)
    bar(s)
    header(s, "COMPETITIVE ADVANTAGE", "Why 700+ Students Chose Mantra Over Everyone Else",
         "11 / " + N,
         "When your future is on the line, the counselor matters more than the brochure.")
    pic(s, "founder.jpg", Inches(0.55), Inches(1.4), Inches(3.0), Inches(3.8))
    txt(s, Inches(0.55), Inches(5.3), Inches(3.0), Inches(0.55),
        "Vinodhini Y.\nFounder & Lead Counselor", 11, True, PRIMARY, PP_ALIGN.CENTER)
    txt(s, Inches(0.55), Inches(5.85), Inches(3.0), Inches(0.45),
        "British Council Certified\nFormer IDP · 10+ Years · Top 17 India",
        8, False, MUTED, PP_ALIGN.CENTER)
    compare_bad = [
        "50 students assigned to one call-center agent",
        "Pushes universities paying highest commission",
        "Generic SOP templates — same essay for everyone",
        "No support after application is submitted",
        "Hidden fees discovered mid-process",
        "30+ countries — shallow expertise in all",
    ]
    compare_good = [
        "ONE dedicated counselor assigned to YOU from day one",
        "University matching based on profile fit, not commission",
        "Custom SOP strategy — your story, expertly told",
        "Full journey: counseling → visa → arrival → alumni network",
        "Transparent pricing — free first counseling session",
        "Deep focus: USA, UK, Australia only — 15+ partner unis",
    ]
    rect(s, Inches(3.8), Inches(1.4), Inches(4.5), Inches(0.45), RGBColor(253, 232, 232))
    txt(s, Inches(3.95), Inches(1.48), Inches(4.2), Inches(0.3),
        "TYPICAL CONSULTANCY", 10, True, RED, PP_ALIGN.CENTER)
    for i, item in enumerate(compare_bad):
        txt(s, Inches(3.95), Inches(1.95 + i * 0.72), Inches(4.3), Inches(0.65),
            f"✗  {item}", 9, False, PRIMARY)
    rect(s, Inches(8.5), Inches(1.4), Inches(4.3), Inches(0.45), PRIMARY)
    txt(s, Inches(8.65), Inches(1.48), Inches(4.0), Inches(0.3),
        "MANTRA GLOBAL EDUCATION", 10, True, ACCENT, PP_ALIGN.CENTER)
    for i, item in enumerate(compare_good):
        txt(s, Inches(8.65), Inches(1.95 + i * 0.72), Inches(4.1), Inches(0.65),
            f"✓  {item}", 9, False, PRIMARY)
    metrics_row = [("700+", "Students\nPlaced"), ("94%", "Partner\nSatisfaction"),
                     ("15+", "Partner\nUniversities"), ("Free", "First\nCounseling")]
    for i, (v, l) in enumerate(metrics_row):
        metric_card(s, Inches(3.8 + i * 2.35), Inches(6.35), Inches(2.15), Inches(0.95), v, l)
    notes(s, "Parents in the room: emphasize British Council certification and transparent process. "
             "Students: emphasize 1:1 counselor and custom SOP.")
    footer(s)

    # ═══ SLIDE 12 — FINANCING ═══
    s = prs.slides.add_slide(B)
    bg(s, WHITE)
    bar(s)
    header(s, "ADDRESSING YOUR #1 CONCERN", "How Families Actually Fund Study Abroad",
         "12 / " + N,
         "We've helped hundreds of middle-class families plan affordable international education.")
    section_card(s, Inches(0.55), Inches(1.4), Inches(6.1), Inches(2.5),
                 "Scholarship Opportunities We Pursue",
                 "Merit scholarships: $5,000–$25,000/year at US universities\n"
                 "Teaching & research assistantships: tuition waiver + stipend\n"
                 "UK university bursaries: £3,000–£10,000 automatic for strong profiles\n"
                 "Australia Awards & university-specific grants\n"
                 "External: Fulbright, Commonwealth, Inlaks (for eligible students)\n\n"
                 "Mantra's role: identify every scholarship your profile qualifies for and "
                 "help craft compelling scholarship essays.")
    section_card(s, Inches(0.55), Inches(4.05), Inches(6.1), Inches(2.7),
                 "Education Loan Partners & Planning",
                 "Indian banks: SBI Global Ed-Vantage, HDFC Credila, Axis Bank\n"
                 "International: Prodigy Finance, MPOWER (no collateral for top unis)\n"
                 "Typical coverage: 80–100% of tuition + living costs\n"
                 "Moratorium: repay only after graduation + 6 months\n\n"
                 "We sit with your family, map total cost, and build a funding plan "
                 "that doesn't create financial stress.")
    section_card(s, Inches(6.85), Inches(1.4), Inches(6.0), Inches(5.35),
                 "Sample Budget Planning (US MS — 2 Years)",
                 "Tuition (avg. public university): ₹25–35L\n"
                 "Living expenses (2 years): ₹16–20L\n"
                 "Total investment: ₹40–55L\n"
                 "Less: Scholarship (avg. Mantra student): ₹3–8L\n"
                 "Less: Part-time work (20 hrs/week): ₹6–10L\n"
                 "Net family investment: ₹25–40L\n\n"
                 "Expected US salary after graduation: ₹80L–1.2Cr/year\n"
                 "ROI break-even period: 2.5–3.5 years\n\n"
                 "The math works — if you plan it with the right guidance.")
    notes(s, "This slide wins parents. If a parent is in the room, make eye contact here. "
             "Say: 'We will build YOUR family's specific budget plan in the free counseling session.'")
    footer(s)

    # ═══ SLIDE 13 — TIMELINE ═══
    s = prs.slides.add_slide(B)
    bg(s, WHITE)
    bar(s)
    header(s, "ACTION PLAN", "Your Study Abroad Timeline — Start Now, Not Later",
         "13 / " + N,
         "Fall 2026 intake deadlines are closer than you think. Here is the roadmap.")
    timeline = [
        ("NOW", "Month 0", "Book free counseling · Profile assessment · Set target country & intake"),
        ("Month 1–2", "Prepare", "IELTS/GRE coaching · Shortlist 8–10 universities · Start SOP draft"),
        ("Month 3–4", "Apply", "Submit applications · Scholarship essays · LOR coordination"),
        ("Month 5–6", "Decide", "Compare admission offers · Accept best fit · Pay deposit"),
        ("Month 7–8", "Visa", "Financial docs · Visa filing · Mock interview prep"),
        ("Month 9+", "Fly", "Pre-departure briefing · Accommodation · Forex · Take off"),
    ]
    for i, (when, phase, action) in enumerate(timeline):
        y = Inches(1.45 + i * 0.88)
        rect(s, Inches(0.55), y, Inches(1.3), Inches(0.72), PRIMARY if i == 0 else LIGHT,
             ACCENT if i == 0 else PRIMARY)
        tf = s.shapes[-1].text_frame
        tf.paragraphs[0].text = when
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = ACCENT if i == 0 else PRIMARY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        txt(s, Inches(2.0), y + Inches(0.04), Inches(1.5), Inches(0.25), phase, 9, True, ACCENT)
        txt(s, Inches(2.0), y + Inches(0.3), Inches(10.5), Inches(0.4), action, 11, False, PRIMARY)
        if i < 5:
            s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.05), y + Inches(0.75), Inches(0.25), Inches(0.15))
    rect(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.42), RGBColor(253, 232, 232))
    txt(s, Inches(0.7), Inches(6.76), Inches(11.9), Inches(0.3),
        "Students who start 12+ months early get better universities, more scholarships, and less stress.",
        10, True, RED, PP_ALIGN.CENTER)
    notes(s, "CREATE URGENCY. 'If you want Fall 2026, your application window is NOW. "
             "Book counseling this week — not next month.'")
    footer(s)

    # ═══ SLIDE 14 — CLOSE / CTA ═══
    s = prs.slides.add_slide(B)
    bg(s, PRIMARY)
    bar(s)
    txt(s, Inches(0.7), Inches(0.45), Inches(5), Inches(0.28),
        "YOUR NEXT STEP", 9, True, ACCENT)
    txt(s, Inches(0.7), Inches(0.75), Inches(10), Inches(0.55),
        "Book Your Free 1-on-1 Counseling Session", 28, True, WHITE, font="Georgia")
    txt(s, Inches(0.7), Inches(1.35), Inches(10), Inches(0.4),
        "Limited slots available this week — first come, first served.", 13, False, GREY_TXT)
    steps = [
        ("1", "Scan & Register",
         "Visit mantraglobaledu.com or meet our desk after this session. "
         "Fill a 2-minute profile form so we come prepared."),
        ("2", "Free Profile Check",
         "30-minute 1-on-1 with a certified counselor. "
         "We assess your academics, goals, budget, and timeline — honestly."),
        ("3", "Get Your Roadmap",
         "Walk away with a personalized study-abroad plan: "
         "country, universities, costs, scholarships, and application timeline."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(2.0 + i * 1.45)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        txt(s, Inches(1.4), y + Inches(0.02), Inches(5.5), Inches(0.3), title, 14, True, WHITE)
        txt(s, Inches(1.4), y + Inches(0.35), Inches(5.8), Inches(0.9), desc, 10, False, GREY_TXT)
    panel = rect(s, Inches(7.5), Inches(1.9), Inches(5.2), Inches(4.5), RGBColor(13, 50, 90))
    txt(s, Inches(7.75), Inches(2.1), Inches(4.7), Inches(0.8),
        "What you get in\nfree counseling:", 16, True, ACCENT, font="Georgia")
    free_items = [
        "Personal profile assessment (₹5,000 value)",
        "Country & university shortlist for YOUR profile",
        "Realistic budget & scholarship assessment",
        "Application timeline with key deadlines",
        "Honest answer to every question you have",
    ]
    bullets(s, Inches(7.75), Inches(3.0), Inches(4.7), Inches(2.5), free_items, 11, GREY_TXT)
    txt(s, Inches(7.75), Inches(5.55), Inches(4.7), Inches(0.35),
        "mantraglobaledu.com", 20, True, ACCENT, PP_ALIGN.CENTER)
    txt(s, Inches(7.75), Inches(5.95), Inches(4.7), Inches(0.3),
        "Coimbatore, Tamil Nadu", 10, False, GREY_TXT, PP_ALIGN.CENTER)
    txt(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
        "Thank you for your time today. Your global future is one conversation away. "
        "Meet us at the desk — or visit mantraglobaledu.com now.",
        12, True, SOFT_GOLD, PP_ALIGN.CENTER)
    notes(s, "CLOSE STRONG. 'Who is ready to book their free session?' "
             "Have signup sheets / QR code ready. Count signups publicly: social proof drives more.")
    footer(s, right="Slide 14 / " + N)

    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
